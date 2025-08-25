from pptx import Presentation

# 프레젠테이션 객체 생성
prs = Presentation()

# 1번 슬라이드 - 제목
slide_layout = prs.slide_layouts[0]  # 제목 슬라이드
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "생활습관 기반 질병 위험도 예측기"
subtitle.text = "20814 서슬빈"

# 2번 슬라이드 - 주제 선정 이유
slide_layout = prs.slide_layouts[1]  # 제목+내용
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "주제 선정 이유"
content.text = (
    "제 진로는 간호, 생명, 식품 분야와 관련이 있습니다.\n"
    "이 웹은 생활습관과 질병의 상관관계를 탐구하면서\n"
    "보건과 영양, 건강 관리의 중요성을 학습할 수 있어\n"
    "진로와 밀접하게 연관된 주제라고 생각했습니다."
)

# 3번 슬라이드 - 빈 슬라이드
slide_layout = prs.slide_layouts[6]  # 완전 빈
prs.slides.add_slide(slide_layout)

# 4번 슬라이드 - 빈 슬라이드
prs.slides.add_slide(slide_layout)

# 저장
prs.save("생활습관_위험도_발표자료.pptx")
print("생활습관_위험도_발표자료.pptx 파일이 생성되었습니다.")
